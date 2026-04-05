"""Wiki-scoped compilation: compile articles from a specific wiki's S3 files into Postgres."""
import asyncio
import json
import logging
import re
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from sqlalchemy import select, delete
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.aws import read_s3_text, get_s3_client
from app.core.config import settings
from app.core.llm import LLMResponse, ask_with_preset, load_prompt
from app.models.wiki import Wiki
from app.models.wiki_article import WikiArticle
from app.models.wiki_file import WikiFile

logger = logging.getLogger(__name__)


class CostTracker:
    """Accumulate token usage and cost across multiple LLM calls."""
    def __init__(self):
        self.prompt_tokens = 0
        self.completion_tokens = 0
        self.total_tokens = 0
        self.cost_usd = 0.0
        self.calls = 0
        self.by_model: dict[str, dict] = {}

    def track(self, resp: LLMResponse):
        self.prompt_tokens += resp.prompt_tokens
        self.completion_tokens += resp.completion_tokens
        self.total_tokens += resp.total_tokens
        self.cost_usd += resp.cost_usd
        self.calls += 1
        if resp.model not in self.by_model:
            self.by_model[resp.model] = {"prompt_tokens": 0, "completion_tokens": 0, "cost_usd": 0.0, "calls": 0}
        m = self.by_model[resp.model]
        m["prompt_tokens"] += resp.prompt_tokens
        m["completion_tokens"] += resp.completion_tokens
        m["cost_usd"] += resp.cost_usd
        m["calls"] += 1

    def to_dict(self) -> dict:
        return {
            "prompt_tokens": self.prompt_tokens,
            "completion_tokens": self.completion_tokens,
            "total_tokens": self.total_tokens,
            "cost_usd": round(self.cost_usd, 4),
            "calls": self.calls,
            "by_model": {
                model: {**data, "cost_usd": round(data["cost_usd"], 4)}
                for model, data in self.by_model.items()
            },
        }

TEXT_TYPES = {"text/markdown", "text/plain"}
DOCX_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_TYPE = "application/pdf"


async def _load_wiki_files(wiki_id: str, db: AsyncSession) -> dict[str, str]:
    """Download and extract text from all clean files in a wiki."""
    result = await db.execute(
        select(WikiFile).where(
            WikiFile.wiki_id == wiki_id,
            WikiFile.status == "clean",
        )
    )
    files = result.scalars().all()
    contents = {}

    for f in files:
        try:
            text = await asyncio.to_thread(_extract_text, f.s3_key, f.content_type)
            if text.strip():
                contents[f.filename] = text
        except Exception as e:
            logger.warning("Failed to extract text from %s: %s", f.filename, e)

    return contents


def _extract_text(s3_key: str, content_type: str) -> str:
    if content_type in TEXT_TYPES:
        return read_s3_text(s3_key)

    s3 = get_s3_client()
    with tempfile.NamedTemporaryFile(suffix=_extension(content_type), delete=True) as tmp:
        s3.download_file(settings.s3_wiki_bucket, s3_key, tmp.name)
        if content_type == PDF_TYPE:
            return _extract_pdf(Path(tmp.name))
        elif content_type == DOCX_TYPE:
            return _extract_docx(Path(tmp.name))
        elif content_type == PPTX_TYPE:
            return _extract_pptx(Path(tmp.name))

    return read_s3_text(s3_key)


def _extension(content_type: str) -> str:
    return {PDF_TYPE: ".pdf", DOCX_TYPE: ".docx", PPTX_TYPE: ".pptx"}.get(content_type, ".txt")


def _extract_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return path.read_text(errors="replace")


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _fit_sources_to_budget(sources: dict[str, str], max_chars: int = 400_000) -> dict[str, str]:
    if not sources:
        return sources
    total = sum(len(v) for v in sources.values())
    if total <= max_chars:
        return sources

    per_source = max_chars // len(sources)
    fitted = {}
    surplus = 0
    long_sources = []

    for name, content in sources.items():
        if len(content) <= per_source:
            fitted[name] = content
            surplus += per_source - len(content)
        else:
            long_sources.append(name)

    if long_sources:
        extra_per = surplus // len(long_sources)
        limit = per_source + extra_per
        for name in long_sources:
            fitted[name] = sources[name][:limit]

    return fitted


def _parse_plan_json(response: str) -> list[dict]:
    try:
        return json.loads(response)
    except json.JSONDecodeError:
        pass
    match = re.search(r"```(?:json)?\s*\n?(.*?)\n?```", response, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass
    match = re.search(r"\[.*\]", response, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass
    raise ValueError(f"Could not parse article plan: {response[:500]}")


# ---------- Compile ----------


async def _get_wiki_models(wiki_id: str, db: AsyncSession) -> dict[str, str]:
    """Get per-wiki model presets."""
    result = await db.execute(select(Wiki).where(Wiki.id == wiki_id))
    wiki = result.scalar_one_or_none()
    if not wiki:
        return {"compile": "groq", "polish": "openai", "qa": "openai"}
    return {
        "compile": wiki.compile_model,
        "polish": wiki.polish_model,
        "qa": wiki.qa_model,
    }


async def compile_wiki_scoped(
    wiki_id: str,
    db: AsyncSession,
    full_rebuild: bool = False,
    on_progress=None,
) -> list[dict]:
    """Three-stage compile pipeline:
    1. Groq (fast/cheap) → bulk write all articles
    2. GPT-5.4 (smart) → polish pass to refine quality
    3. Save to Postgres

    Models are configurable per-wiki.
    on_progress: async callback(current, total, message).
    Returns list of {slug, title, category}.
    """
    async def progress(current, total, message):
        if on_progress:
            await on_progress(current, total, message)

    models = await _get_wiki_models(wiki_id, db)
    costs = CostTracker()

    await progress(0, 0, "Loading source files...")
    raw_full = await _load_wiki_files(wiki_id, db)
    if not raw_full:
        return [], costs

    await progress(0, 0, f"Loaded {len(raw_full)} files. Planning with {models['compile']}...")
    raw_summaries = {name: content[:4000] for name, content in raw_full.items()}

    wiki_state = {}
    if not full_rebuild:
        wiki_state = await _load_existing_articles(wiki_id, db)

    plan, plan_resp = await asyncio.to_thread(
        _plan_articles_sync, raw_summaries, wiki_state, models["compile"]
    )
    costs.track(plan_resp)
    total = len(plan)

    await progress(0, total, f"Planned {total} articles. Drafting with {models['compile']}...")

    if full_rebuild:
        await db.execute(
            delete(WikiArticle).where(WikiArticle.wiki_id == wiki_id)
        )
        await db.flush()

    # Stage 1: Parallel draft with fast/cheap model
    sem = asyncio.Semaphore(10)  # max 10 concurrent LLM calls
    draft_count = 0

    async def draft_one(article):
        nonlocal draft_count
        async with sem:
            try:
                content, resp = await asyncio.to_thread(
                    _write_article_sync, article, raw_full, models["compile"]
                )
                costs.track(resp)
                draft_count += 1
                await progress(draft_count, total, f"Drafted \"{article['title']}\" ({models['compile']})")
                return {**article, "_content": content}
            except Exception as e:
                logger.error("Failed to draft article %s: %s", article.get("slug"), e)
                return None

    draft_results = await asyncio.gather(*[draft_one(a) for a in plan])
    drafts = [d for d in draft_results if d is not None]

    # Stage 2: Parallel polish with smart model
    written = []
    now = datetime.now(timezone.utc)
    polish_count = 0

    async def polish_one(article):
        nonlocal polish_count
        async with sem:
            try:
                polished, resp = await asyncio.to_thread(
                    _polish_article_sync, article["_content"], article, models["polish"]
                )
                costs.track(resp)
                polish_count += 1
                await progress(polish_count, total, f"Polished \"{article['title']}\" ({models['polish']})")
                return {**article, "_polished": polished}
            except Exception as e:
                logger.error("Failed to polish article %s, will use draft: %s", article.get("slug"), e)
                return {**article, "_polished": article["_content"]}

    polish_results = await asyncio.gather(*[polish_one(a) for a in drafts])

    # Stage 3: Save all to DB (sequential — DB writes are fast)
    for article in polish_results:
        existing = await db.execute(
            select(WikiArticle).where(
                WikiArticle.wiki_id == wiki_id,
                WikiArticle.slug == article["slug"],
            )
        )
        existing_article = existing.scalar_one_or_none()

        if existing_article:
            existing_article.title = article["title"]
            existing_article.category = article["category"]
            existing_article.content = article["_polished"]
            existing_article.sources = article.get("sources", [])
            existing_article.tags = article.get("tags", [])
            existing_article.updated_at = now
        else:
            db.add(WikiArticle(
                wiki_id=wiki_id,
                slug=article["slug"],
                title=article["title"],
                category=article["category"],
                content=article["_polished"],
                sources=article.get("sources", []),
                tags=article.get("tags", []),
            ))

        written.append({
            "slug": article["slug"],
            "title": article["title"],
            "category": article["category"],
        })

    await progress(total, total, "Saving to database...")
    await db.commit()
    return written, costs


async def _load_existing_articles(wiki_id: str, db: AsyncSession) -> dict[str, str]:
    """Load existing compiled articles from Postgres (first 2000 chars each)."""
    result = await db.execute(
        select(WikiArticle).where(WikiArticle.wiki_id == wiki_id)
    )
    articles = result.scalars().all()
    return {
        f"{a.category}/{a.slug}.md": a.content[:2000]
        for a in articles
    }


def _plan_articles_sync(
    raw_summaries: dict[str, str], wiki_state: dict[str, str], model_preset: str = "groq"
) -> list[dict]:
    system = load_prompt("compilation_plan_system")

    parts = ["## Raw Source Documents\n"]
    for name, content in raw_summaries.items():
        parts.append(f"### {name}\n{content}\n")

    if wiki_state:
        parts.append("\n## Existing Wiki Articles\n")
        for name, content in wiki_state.items():
            parts.append(f"### {name}\n{content}\n")
    else:
        parts.append("\n## Existing Wiki Articles\nNone — this is a fresh build.\n")

    parts.append("\nNow plan the wiki articles. Return ONLY the JSON array.")
    result = ask_with_preset(model_preset, system, "\n".join(parts), max_tokens=16384)
    return _parse_plan_json(result.text), result


def _write_article_sync(
    article: dict, raw_full: dict[str, str], model_preset: str = "groq"
) -> str:
    system = load_prompt("compilation_write_system")

    relevant_files = {}
    for source_name in article.get("sources", []):
        if source_name in raw_full:
            relevant_files[source_name] = raw_full[source_name]

    relevant_files = _fit_sources_to_budget(relevant_files, max_chars=400_000)

    user_message = (
        f"## Article to write\n"
        f"Title: {article['title']}\n"
        f"Category: {article['category']}\n"
        f"Description: {article.get('description', '')}\n"
        f"Related articles (use [[wikilinks]] to these): {', '.join(article.get('tags', []))}\n"
        f"\nWrite this article now."
    )

    result = ask_with_preset(model_preset, system, user_message, relevant_files, max_tokens=8192)
    return result.text, result


POLISH_SYSTEM = """\
You are an expert medical editor reviewing a wiki article for a graduate medical \
education knowledge base. The article was drafted by a fast AI model and needs refinement.

Your job:
1. Fix any medical inaccuracies or imprecise terminology
2. Add important caveats and qualifications where needed
3. Strengthen [[wikilinks]] to other articles — add links where related concepts are mentioned
4. Improve clarity and readability without changing the structure
5. Ensure citations to source documents are accurate
6. Remove any hallucinated facts not supported by the source material

Return the FULL improved article in markdown. Keep the same structure and length — \
this is a polish, not a rewrite. If the article is already good, return it with minimal changes."""


def _polish_article_sync(
    draft: str, article: dict, model_preset: str = "openai"
) -> str:
    user_message = (
        f"## Article to polish\n"
        f"Title: {article['title']}\n"
        f"Category: {article['category']}\n\n"
        f"## Draft\n{draft}\n\n"
        f"Polish this article now. Return the full improved markdown."
    )

    result = ask_with_preset(model_preset, POLISH_SYSTEM, user_message, max_tokens=8192)
    return result.text, result


# ---------- Lint ----------


async def lint_wiki_scoped(wiki_id: str, db: AsyncSession, on_progress=None) -> str:
    """Lint a specific wiki's compiled articles. Returns the report."""
    async def progress(current, total, message):
        if on_progress:
            await on_progress(current, total, message)

    await progress(0, 0, "Loading source files...")
    source_files = await _load_wiki_files(wiki_id, db)

    result = await db.execute(
        select(WikiArticle).where(WikiArticle.wiki_id == wiki_id)
    )
    articles = result.scalars().all()

    if not articles:
        return "No compiled articles found. Run compile first."

    await progress(0, 1, f"Auditing {len(articles)} articles...")
    compiled = {
        f"{a.category}/{a.slug}.md": a.content
        for a in articles
    }

    models = await _get_wiki_models(wiki_id, db)
    costs = CostTracker()
    report, resp = await asyncio.to_thread(
        _lint_sync, source_files, compiled, models["polish"]
    )
    costs.track(resp)
    await progress(1, 1, "Lint complete")
    return report, costs


def _lint_sync(
    source_files: dict[str, str], compiled: dict[str, str], model_preset: str = "openai"
) -> str:
    system = (
        "You are a knowledge base auditor. Review the wiki articles and report on:\n"
        "1. **Contradictions** — conflicting facts across articles\n"
        "2. **Broken wikilinks** — [[links]] referencing non-existent articles\n"
        "3. **Coverage gaps** — source files not represented in articles\n"
        "4. **Missing cross-references** — related articles not linking to each other\n"
        "5. **Thin articles** — too short or lacking concrete details\n"
        "6. **Suggested new articles** — concepts that deserve their own article\n\n"
        "Be specific: name the article, quote text, suggest fixes.\n"
        "End with an overall health score (1-10) and top 3 priorities."
    )

    file_contents = {}
    for name, content in compiled.items():
        file_contents[name] = content[:4000]

    raw_listing = "\n".join(f"- {name}" for name in source_files.keys())
    file_contents["_source_file_list.txt"] = raw_listing

    user_msg = (
        f"Review this wiki. It has {len(compiled)} compiled articles "
        f"from {len(source_files)} source files.\n\n"
        f"Produce a detailed quality report."
    )

    result = ask_with_preset(model_preset, system, user_msg, file_contents, max_tokens=8192)
    return result.text, result


# ---------- Q&A ----------


async def ask_wiki_scoped(wiki_id: str, question: str, db: AsyncSession) -> str:
    """Answer a question using a specific wiki's compiled articles + wiki's QA model."""
    models = await _get_wiki_models(wiki_id, db)

    result = await db.execute(
        select(WikiArticle).where(WikiArticle.wiki_id == wiki_id)
    )
    articles = result.scalars().all()

    if not articles:
        return "No compiled articles found. Run compile first."

    file_contents = {}
    for a in articles:
        file_contents[f"{a.category}/{a.slug}.md"] = a.content[:6000]

    answer, resp = await asyncio.to_thread(
        _ask_sync, question, file_contents, models["qa"]
    )
    return answer, resp


def _ask_sync(question: str, file_contents: dict[str, str], model_preset: str = "openai") -> str:
    system = (
        "You are a research assistant with access to a wiki knowledge base. "
        "Answer thoroughly using the provided articles. "
        "Cite articles with [[wikilinks]]. "
        "Say clearly if the information is insufficient to answer."
    )
    result = ask_with_preset(model_preset, system, question, file_contents, max_tokens=4096)
    return result.text, result
