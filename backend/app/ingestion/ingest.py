import re
import shutil
import subprocess
from pathlib import Path

import frontmatter
import html2text
import httpx
from bs4 import BeautifulSoup
from rich import print as rprint

from app.core.config import settings


def slugify(text: str) -> str:
    """Filesystem-safe slug: lowercase, no special chars, max 80 chars."""
    slug = text.lower().strip()
    slug = re.sub(r"[^\w\s-]", "", slug)
    slug = re.sub(r"[\s_-]+", "-", slug)
    slug = slug.strip("-")
    return slug[:80]


def _fetch_url(url: str) -> str:
    """Fetch URL content. Try httpx first, fall back to urllib for sites that block it."""
    headers = {"User-Agent": "Mozilla/5.0 (X11; Linux x86_64; rv:128.0) Gecko/20100101 Firefox/128.0"}
    try:
        response = httpx.get(url, follow_redirects=True, timeout=30, headers=headers)
        response.raise_for_status()
        return response.text
    except httpx.HTTPStatusError:
        import urllib.request

        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=30) as resp:
            return resp.read().decode(resp.headers.get_content_charset() or "utf-8")


def ingest_url(url: str) -> Path:
    """Fetch URL, convert HTML to markdown, write to raw/."""
    html = _fetch_url(url)

    soup = BeautifulSoup(html, "html.parser")
    title_tag = soup.find("title")
    title = title_tag.get_text(strip=True) if title_tag else url

    converter = html2text.HTML2Text()
    converter.ignore_links = False
    converter.ignore_images = True
    converter.body_width = 0
    md_content = converter.handle(html)

    slug = slugify(title)
    path = settings.raw_dir / f"{slug}.md"

    metadata = {"title": title, "source": url, "type": "web"}
    post = frontmatter.Post(md_content, **metadata)
    path.write_text(frontmatter.dumps(post) + "\n")

    rprint(f"[green]Ingested URL:[/green] {title} → {path.name}")
    return path


def ingest_pdf(pdf_path: Path) -> Path:
    """Convert PDF to markdown. Try marker subprocess first, fall back to PyPDF2."""
    title = pdf_path.stem
    slug = slugify(title)
    out_path = settings.raw_dir / f"{slug}.md"

    md_content = _try_marker(pdf_path)
    if md_content is None:
        md_content = _fallback_pypdf2(pdf_path)

    metadata = {"title": title, "source": str(pdf_path), "type": "pdf"}
    post = frontmatter.Post(md_content, **metadata)
    out_path.write_text(frontmatter.dumps(post) + "\n")

    rprint(f"[green]Ingested PDF:[/green] {title} → {out_path.name}")
    return out_path


def _try_marker(pdf_path: Path) -> str | None:
    """Try converting PDF with marker_single subprocess."""
    try:
        result = subprocess.run(
            ["marker_single", str(pdf_path), str(settings.raw_dir), "--max_pages", "100"],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0:
            # marker_single writes output to a subdirectory; find the .md file
            marker_out = settings.raw_dir / pdf_path.stem
            md_files = list(marker_out.glob("*.md"))
            if md_files:
                content = md_files[0].read_text()
                # Clean up marker output directory
                shutil.rmtree(marker_out, ignore_errors=True)
                return content
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None


def _fallback_pypdf2(pdf_path: Path) -> str:
    """Extract text from PDF using PyPDF2."""
    from PyPDF2 import PdfReader

    reader = PdfReader(str(pdf_path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages.append(f"## Page {i + 1}\n\n{text}")
    return "\n\n".join(pages) if pages else "_No extractable text found._"


def ingest_markdown(md_path: Path) -> Path:
    """Copy markdown to raw/, inject frontmatter if missing."""
    content = md_path.read_text()
    slug = slugify(md_path.stem)
    out_path = settings.raw_dir / f"{slug}.md"

    try:
        post = frontmatter.loads(content)
        if not post.metadata:
            raise ValueError("empty frontmatter")
    except Exception:
        post = frontmatter.Post(content, title=md_path.stem, source=str(md_path), type="markdown")

    out_path.write_text(frontmatter.dumps(post) + "\n")
    rprint(f"[green]Ingested markdown:[/green] {md_path.name} → {out_path.name}")
    return out_path


def ingest_docx(docx_path: Path) -> Path:
    """Convert DOCX to markdown, write to raw/."""
    from docx import Document

    doc = Document(str(docx_path))
    title = docx_path.stem
    slug = slugify(title)
    out_path = settings.raw_dir / f"{slug}.md"

    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        elif "heading" in style:
            lines.append(f"#### {text}")
        else:
            lines.append(text)

    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Add header separator after first row
            header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
            rows.insert(1, header_sep)
            lines.append("\n".join(rows))

    md_content = "\n\n".join(lines)
    metadata = {"title": title, "source": str(docx_path), "type": "docx"}
    post = frontmatter.Post(md_content, **metadata)
    out_path.write_text(frontmatter.dumps(post) + "\n")

    rprint(f"[green]Ingested DOCX:[/green] {title} → {out_path.name}")
    return out_path


def ingest_pptx(pptx_path: Path) -> Path:
    """Convert PPTX to markdown, write to raw/."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    title = pptx_path.stem
    slug = slugify(title)
    out_path = settings.raw_dir / f"{slug}.md"

    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_body = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # First text on a slide is usually the title
                    if not slide_title and shape.shape_type is not None:
                        slide_title = text
                    else:
                        slide_body.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header_sep = "| " + " | ".join("---" for _ in shape.table.rows[0].cells) + " |"
                    rows.insert(1, header_sep)
                    slide_body.append("\n".join(rows))

        heading = slide_title or f"Slide {i}"
        lines.append(f"## {heading}")
        if slide_body:
            lines.append("\n".join(slide_body))

    md_content = "\n\n".join(lines)
    metadata = {"title": title, "source": str(pptx_path), "type": "pptx"}
    post = frontmatter.Post(md_content, **metadata)
    out_path.write_text(frontmatter.dumps(post) + "\n")

    rprint(f"[green]Ingested PPTX:[/green] {title} → {out_path.name}")
    return out_path


SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".md": "markdown",
    ".markdown": "markdown",
    ".txt": "markdown",
    ".docx": "docx",
    ".pptx": "pptx",
}


def ingest_directory(dir_path: Path) -> list[Path]:
    """Iterate directory and dispatch to appropriate handler."""
    results = []
    for item in sorted(dir_path.iterdir()):
        if item.is_dir():
            results.extend(ingest_directory(item))
        elif item.suffix in SUPPORTED_EXTENSIONS:
            results.append(ingest(str(item)))
        else:
            rprint(f"[yellow]Skipping unsupported file:[/yellow] {item.name}")
    return results


def ingest(source: str) -> Path | list[Path]:
    """Auto-detect source type and dispatch."""
    if source.startswith(("http://", "https://")):
        return ingest_url(source)

    path = Path(source).resolve()
    if not path.exists():
        raise FileNotFoundError(f"Source not found: {source}")

    if path.is_dir():
        return ingest_directory(path)
    elif path.suffix == ".pdf":
        return ingest_pdf(path)
    elif path.suffix in (".md", ".markdown", ".txt"):
        return ingest_markdown(path)
    elif path.suffix == ".docx":
        return ingest_docx(path)
    elif path.suffix == ".pptx":
        return ingest_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {path.suffix}")
